# -*- coding: utf-8 -*-
"""
优品汇电商全链路经营分析 - 面试答辩PPT生成脚本
生成16:9宽屏PPT，包含18页完整分析汇报
"""

import os
import sys
import json
import pandas as pd
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Cm, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ============ 路径配置 ============
BASE_DIR = Path(__file__).resolve().parent.parent
DATA_DIR = BASE_DIR / "data" / "processed"
PLOTS_DIR = DATA_DIR / "plots"
OUTPUT_PATH = BASE_DIR / "report" / "优品汇电商分析_面试汇报.pptx"

# ============ 配色方案 ============
COLOR_TITLE = RGBColor(0, 51, 102)       # 深蓝标题
COLOR_SUBTITLE = RGBColor(102, 102, 102) # 灰色副标题
COLOR_WHITE = RGBColor(255, 255, 255)
COLOR_BLACK = RGBColor(0, 0, 0)
COLOR_ACCENT = RGBColor(0, 102, 178)     # 强调蓝
COLOR_HIGHLIGHT = RGBColor(192, 57, 43)  # 强调红
COLOR_BG_DARK = RGBColor(0, 51, 102)     # 深蓝背景

FONT_NAME = "Microsoft YaHei"
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


# ============ 数据加载 ============
def load_data():
    """加载所有分析结果数据"""
    data = {}
    # 模型对比
    data['model'] = pd.read_csv(DATA_DIR / "model_comparison.csv")
    # 因果推断
    with open(DATA_DIR / "causal_results.json", "r", encoding="utf-8") as f:
        data['causal'] = json.load(f)
    # 月度KPI
    data['kpi'] = pd.read_csv(DATA_DIR / "monthly_kpi.csv")
    # 品类表现
    data['category'] = pd.read_csv(DATA_DIR / "category_performance.csv", encoding="utf-8")
    return data


# ============ 辅助函数 ============
def set_font(run, size=18, bold=False, color=COLOR_BLACK, name=FONT_NAME):
    """设置字体属性"""
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = name


def add_title_bar(slide, title_text, subtitle_text=None):
    """添加页面顶部标题栏"""
    # 标题背景色块
    left, top, width, height = Inches(0), Inches(0), SLIDE_WIDTH, Inches(1.2)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_BG_DARK
    shape.line.fill.background()

    # 标题文字
    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(0.15), Inches(10), Inches(0.9))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title_text
    set_font(run, size=28, bold=True, color=COLOR_WHITE)

    # 副标题
    if subtitle_text:
        txBox2 = slide.shapes.add_textbox(Inches(0.6), Inches(1.35), Inches(10), Inches(0.5))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        run2 = p2.add_run()
        run2.text = subtitle_text
        set_font(run2, size=14, color=COLOR_SUBTITLE)


def add_bullet_points(slide, bullets, left=0.8, top=2.0, width=11.5, height=4.5, font_size=18):
    """添加项目符号列表"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(12)
        run = p.add_run()
        run.text = f"▸ {bullet}"
        set_font(run, size=font_size, color=COLOR_BLACK)


def add_big_number(slide, number_text, desc_text, left=1.0, top=2.5):
    """添加大数字突出显示"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(5), Inches(1.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = number_text
    set_font(run, size=44, bold=True, color=COLOR_ACCENT)
    p2 = tf.add_paragraph()
    run2 = p2.add_run()
    run2.text = desc_text
    set_font(run2, size=16, color=COLOR_SUBTITLE)


def add_image_centered(slide, img_path, top=2.0, max_width=10.0, max_height=5.0):
    """居中添加图片"""
    if not Path(img_path).exists():
        return
    from PIL import Image
    try:
        img = Image.open(img_path)
        img_w, img_h = img.size
    except Exception:
        img_w, img_h = 800, 600

    # 计算合适的缩放
    ratio = min(max_width / (img_w / 96), max_height / (img_h / 96))
    display_w = (img_w / 96) * ratio
    display_h = (img_h / 96) * ratio
    left = (13.333 - display_w) / 2
    slide.shapes.add_picture(str(img_path), Inches(left), Inches(top),
                             Inches(display_w), Inches(display_h))


def add_image_simple(slide, img_path, left, top, width, height):
    """简单添加图片，指定位置和大小"""
    if Path(img_path).exists():
        slide.shapes.add_picture(str(img_path), Inches(left), Inches(top),
                                 Inches(width), Inches(height))


# ============ 各页生成函数 ============
def slide_cover(prs):
    """第1页：封面"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
    # 深蓝全屏背景
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLOR_BG_DARK
    bg.line.fill.background()

    # 主标题
    txBox = slide.shapes.add_textbox(Inches(1.5), Inches(2.2), Inches(10), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "优品汇电商全链路经营分析"
    set_font(run, size=40, bold=True, color=COLOR_WHITE)

    # 副标题
    txBox2 = slide.shapes.add_textbox(Inches(1.5), Inches(3.8), Inches(10), Inches(1.0))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = "用户运营与流失预测全链路分析"
    set_font(run2, size=24, color=RGBColor(173, 216, 230))

    # 日期和信息
    txBox3 = slide.shapes.add_textbox(Inches(1.5), Inches(5.2), Inches(10), Inches(1.0))
    tf3 = txBox3.text_frame
    p3 = tf3.paragraphs[0]
    p3.alignment = PP_ALIGN.CENTER
    run3 = p3.add_run()
    run3.text = "2026年6月  |  数据分析岗位面试汇报"
    set_font(run3, size=16, color=RGBColor(200, 200, 200))


def slide_toc(prs):
    """第2页：目录"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "目录", "分析框架概览")
    items = [
        "一、项目背景与分析目标",
        "二、技术架构与数据概览",
        "三、RFM用户分层与流失定义",
        "四、流失预测模型",
        "五、因果推断分析",
        "六、NLP评论挖掘",
        "七、经营核心发现与业务建议",
        "八、项目亮点与Q&A"
    ]
    add_bullet_points(slide, items, top=1.8, font_size=20)


def slide_background(prs):
    """第3页：项目背景"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "项目背景", "优品汇企业简介与业务痛点")
    bullets = [
        "优品汇：中国领先综合电商平台，覆盖美妆、数码、家纺、运动等多品类",
        "2016-2018年累计10万+订单，3万+注册用户",
        "核心痛点①：用户复购率低，流失率居高不下",
        "核心痛点②：运营缺乏精细化分层，营销资源浪费严重",
        "核心痛点③：缺少科学的因果归因，难以评估策略效果"
    ]
    add_bullet_points(slide, bullets, top=1.8)


def slide_objectives(prs):
    """第4页：分析目标"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "分析目标", "四大核心分析模块")
    # 四个模块方框
    modules = [
        ("经营分析", "GMV趋势、品类表现\n区域差异、支付行为"),
        ("用户分层", "RFM模型、K-Means\n高价值用户识别"),
        ("流失预测", "动态定义、SMOTE\nLR/RF/XGB对比"),
        ("因果推断", "PSM倾向评分匹配\nATT效应估计")
    ]
    for i, (title, desc) in enumerate(modules):
        left = 0.8 + i * 3.2
        top = 2.2
        # 方框
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                       Inches(left), Inches(top), Inches(2.8), Inches(3.5))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(240, 248, 255)
        shape.line.color.rgb = COLOR_ACCENT
        # 模块标题
        txBox = slide.shapes.add_textbox(Inches(left + 0.2), Inches(top + 0.3),
                                         Inches(2.4), Inches(0.6))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = title
        set_font(run, size=20, bold=True, color=COLOR_ACCENT)
        # 模块描述
        txBox2 = slide.shapes.add_textbox(Inches(left + 0.2), Inches(top + 1.2),
                                          Inches(2.4), Inches(2.0))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.alignment = PP_ALIGN.CENTER
        run2 = p2.add_run()
        run2.text = desc
        set_font(run2, size=14, color=COLOR_BLACK)


def slide_architecture(prs):
    """第5页：技术架构"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "技术架构", "MySQL + Python + Docker 全链路自动化")
    bullets = [
        "数据层：MySQL 8.0 / SQLite 数据存储，支持SQL直接查询",
        "计算层：Python (Pandas + Scikit-learn + XGBoost + CausalInference)",
        "可视化：Matplotlib + Seaborn + WordCloud，自动生成图表",
        "部署层：Docker一键部署，docker-compose编排，完整复现",
        "流程管理：main.py统一调度 → 数据清洗 → 特征工程 → 建模 → 输出报告"
    ]
    add_bullet_points(slide, bullets, top=1.8, font_size=17)


def slide_data_overview(prs):
    """第6页：数据概览"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "数据概览", "数据规模与质量")
    # 大数字展示
    numbers = [
        ("9", "张数据表"),
        ("96,478", "笔订单"),
        ("30,000+", "位用户"),
        ("2016-2018", "时间跨度"),
    ]
    for i, (num, desc) in enumerate(numbers):
        left = 0.8 + i * 3.2
        txBox = slide.shapes.add_textbox(Inches(left), Inches(2.0), Inches(2.8), Inches(2.0))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = num
        set_font(run, size=40, bold=True, color=COLOR_ACCENT)
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        run2 = p2.add_run()
        run2.text = desc
        set_font(run2, size=16, color=COLOR_SUBTITLE)

    bullets = [
        "核心表：orders / order_items / customers / payments / reviews / products",
        "数据清洗：去除异常状态订单，补全缺失品类，时间字段标准化",
        "特征工程：构建50+维度用户特征（RFM、支付偏好、品类偏好、评分行为等）"
    ]
    add_bullet_points(slide, bullets, top=4.2, font_size=16)


def slide_rfm(prs):
    """第7页：RFM用户分层"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "RFM用户分层", "基于购买行为的用户价值分层")
    bullets = [
        "R (Recency)：最近一次购买距今天数 → 活跃度指标",
        "F (Frequency)：购买频次 → 忠诚度指标",
        "M (Monetary)：累计消费金额 → 价值贡献指标",
        "分层方法：各维度按中位数二分 → 2³ = 8个用户群",
        "应用场景：高价值用户维护、沉睡用户唤醒、流失预警"
    ]
    add_bullet_points(slide, bullets, top=1.8)

    # 补充说明
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(5.5), Inches(11), Inches(1.0))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "💡 RFM分层结果作为流失预测模型的关键输入特征"
    set_font(run, size=15, bold=True, color=COLOR_HIGHLIGHT)


def slide_churn_definition(prs):
    """第8页：流失定义创新"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "流失定义创新", "动态阈值 vs 传统固定阈值")
    bullets = [
        "传统方法：固定90天未购买 = 流失（一刀切，忽略品类差异）",
        "本项目创新：购买间隔 P75 × 2 作为动态流失阈值",
        "优势①：适应不同用户购买节奏，高频用户更快识别流失",
        "优势②：基于数据分布而非经验值，更具统计合理性",
        "优势③：可随业务增长动态调整，无需人工重新设定"
    ]
    add_bullet_points(slide, bullets, top=1.8)

    # 对比表格
    txBox = slide.shapes.add_textbox(Inches(1.0), Inches(5.3), Inches(11), Inches(1.2))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "对比：传统90天固定 → 召回率偏低 | P75×2动态 → 精准捕获早期流失信号"
    set_font(run, size=15, bold=True, color=COLOR_ACCENT)


def slide_model_design(prs):
    """第9页：预测模型设计"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "流失预测模型设计", "多模型对比 + 类不平衡处理")
    bullets = [
        "候选模型：Logistic Regression / Random Forest / XGBoost",
        "类不平衡处理：SMOTE过采样，平衡正负样本比例",
        "评估策略：5折交叉验证，防止过拟合",
        "评价指标：AUC-ROC（主要）、F1-Score、Recall（流失类）",
        "特征选择：基于RFM + 支付行为 + 品类偏好 + 评分特征"
    ]
    add_bullet_points(slide, bullets, top=1.8)


def slide_model_results(prs, data):
    """第10页：模型结果"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "模型预测结果", "三模型性能对比")

    model_df = data['model']
    # 关键指标
    metrics_text = []
    for _, row in model_df.iterrows():
        name = row['model']
        auc = row['test_auc_roc']
        f1 = row['test_f1']
        recall = row['test_recall']
        metrics_text.append(f"{name}：AUC={auc:.4f}  F1={f1:.4f}  Recall={recall:.4f}")

    add_bullet_points(slide, metrics_text, top=1.6, font_size=17, height=1.8)

    # ROC曲线图
    roc_path = PLOTS_DIR / "roc_comparison.png"
    add_image_simple(slide, str(roc_path), left=2.5, top=3.4, width=8.0, height=3.8)


def slide_feature_importance(prs):
    """第11页：特征重要性"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "特征重要性分析", "XGBoost模型Top特征")

    img_path = PLOTS_DIR / "feature_importance.png"
    add_image_simple(slide, str(img_path), left=1.5, top=1.5, width=10.0, height=5.5)


def slide_causal_design(prs):
    """第12页：因果推断设计"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "因果推断设计", "倾向评分匹配（PSM）方法")
    bullets = [
        "研究问题：物流延迟是否导致用户流失率上升？",
        "方法：Propensity Score Matching (倾向评分匹配)",
        "处理组：经历过物流延迟的用户 (N=6,353)",
        "对照组：未经历物流延迟的用户 (N=86,402)",
        "协变量：消费金额、评分、品类偏好、支付方式",
        "匹配策略：最近邻匹配 + Caliper=0.05，匹配率97.8%"
    ]
    add_bullet_points(slide, bullets, top=1.8, font_size=17)


def slide_causal_results(prs, data):
    """第13页：因果推断结果"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "因果推断结果", "ATT效应估计")

    causal = data['causal']
    att = causal['att']['att']
    ci_lower = causal['att']['ci_lower']
    ci_upper = causal['att']['ci_upper']
    treat_mean = causal['att']['treatment_mean']
    ctrl_mean = causal['att']['control_mean']

    # 大数字
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(5.5), Inches(2.0))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = f"ATT = {att:.4f} (+{att*100:.2f}%)"
    set_font(run, size=36, bold=True, color=COLOR_HIGHLIGHT)
    p2 = tf.add_paragraph()
    run2 = p2.add_run()
    run2.text = f"95% CI: [{ci_lower:.4f}, {ci_upper:.4f}]"
    set_font(run2, size=18, color=COLOR_SUBTITLE)
    p3 = tf.add_paragraph()
    p3.space_before = Pt(12)
    run3 = p3.add_run()
    run3.text = f"处理组流失率: {treat_mean*100:.2f}% | 对照组: {ctrl_mean*100:.2f}%"
    set_font(run3, size=16, color=COLOR_BLACK)

    # Love plot
    img_path = PLOTS_DIR / "love_plot.png"
    add_image_simple(slide, str(img_path), left=6.5, top=1.5, width=6.0, height=5.2)

    # 结论
    txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(5.5), Inches(5.5), Inches(1.5))
    tf2 = txBox2.text_frame
    p4 = tf2.paragraphs[0]
    run4 = p4.add_run()
    run4.text = "结论：物流延迟显著提升用户流失概率约2.06个百分点"
    set_font(run4, size=15, bold=True, color=COLOR_ACCENT)


def slide_nlp(prs):
    """第14页：NLP评论分析"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "NLP评论挖掘", "用户差评关键词与维度分析")

    # 词云
    wc_path = PLOTS_DIR / "wordcloud_low_score.png"
    add_image_simple(slide, str(wc_path), left=0.3, top=1.5, width=6.2, height=5.5)

    # 维度分析
    dim_path = PLOTS_DIR / "review_dimension_analysis.png"
    add_image_simple(slide, str(dim_path), left=6.8, top=1.5, width=6.2, height=5.5)


def slide_business_findings(prs, data):
    """第15页：经营核心发现"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "经营核心发现", "GMV趋势、品类与区域分析")

    kpi = data['kpi']
    total_gmv = kpi['gmv'].sum()
    total_orders = kpi['order_count'].sum()
    avg_aov = kpi['aov'].mean()
    peak_month = kpi.loc[kpi['gmv'].idxmax()]

    cat = data['category']
    top3 = cat.head(3)['product_category_name'].tolist()

    bullets = [
        f"累计GMV：¥{total_gmv/10000:.1f}万，总订单：{total_orders:,}笔",
        f"平均客单价：¥{avg_aov:.1f}",
        f"峰值月份：{peak_month['month']}，GMV达¥{peak_month['gmv']/10000:.1f}万",
        f"TOP3品类：{top3[0]}、{top3[1]}、{top3[2]}",
        f"GMV整体呈增长趋势，2017年Q2-Q3进入高速增长期"
    ]
    add_bullet_points(slide, bullets, top=1.8, font_size=17)

    # GMV趋势图
    gmv_path = PLOTS_DIR / "business_gmv_trend.png"
    add_image_simple(slide, str(gmv_path), left=3.5, top=4.0, width=7.0, height=3.2)


def slide_recommendations(prs):
    """第16页：业务建议"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "业务建议", "基于数据分析的可落地策略")
    bullets = [
        "① 物流优化：针对延迟率高的区域/品类优先改善配送，预计降低流失率2%+",
        "② 精准营销：基于RFM分层，对高价值沉睡用户定向发放复购激励券",
        "③ 差评预警：NLP实时监测差评关键词，触发客服主动介入挽回",
        "④ 品类运营：加大美妆个护、运动户外等高毛利品类曝光与促销力度",
        "⑤ 流失干预：模型上线后对高流失概率用户（>0.7）自动触发挽留流程"
    ]
    add_bullet_points(slide, bullets, top=1.8, font_size=17)


def slide_highlights(prs):
    """第17页：项目亮点"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "项目亮点", "三大技术创新点")

    highlights = [
        ("🔬 动态流失定义", "基于P75×2的数据驱动流失阈值，替代传统90天固定规则，\n更精准适配不同消费频率用户群体"),
        ("📊 因果推断应用", "引入PSM方法量化物流延迟对流失的因果效应，\n超越相关性分析，为业务干预提供科学依据"),
        ("🐳 全链路自动化", "MySQL → Python → Docker 一键复现，\n从数据清洗到模型训练全流程自动化，可直接部署"),
    ]
    for i, (title, desc) in enumerate(highlights):
        top = 1.7 + i * 1.9
        # 标题
        txBox = slide.shapes.add_textbox(Inches(1.0), Inches(top), Inches(11), Inches(0.6))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = title
        set_font(run, size=22, bold=True, color=COLOR_ACCENT)
        # 描述
        txBox2 = slide.shapes.add_textbox(Inches(1.5), Inches(top + 0.6), Inches(10.5), Inches(1.2))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        run2 = p2.add_run()
        run2.text = desc
        set_font(run2, size=16, color=COLOR_BLACK)


def slide_qa(prs):
    """第18页：Q&A"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # 深蓝全屏背景
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLOR_BG_DARK
    bg.line.fill.background()

    # Q&A
    txBox = slide.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(10), Inches(1.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Q & A"
    set_font(run, size=54, bold=True, color=COLOR_WHITE)

    # 致谢
    txBox2 = slide.shapes.add_textbox(Inches(1.5), Inches(4.5), Inches(10), Inches(1.0))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = "感谢您的聆听与指导"
    set_font(run2, size=24, color=RGBColor(173, 216, 230))

    txBox3 = slide.shapes.add_textbox(Inches(1.5), Inches(5.5), Inches(10), Inches(1.0))
    tf3 = txBox3.text_frame
    p3 = tf3.paragraphs[0]
    p3.alignment = PP_ALIGN.CENTER
    run3 = p3.add_run()
    run3.text = "优品汇电商全链路经营分析项目"
    set_font(run3, size=16, color=RGBColor(200, 200, 200))


# ============ 主函数 ============
def main():
    import io, sys
    sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')

    print("=" * 60)
    print("PPT Generation Start")
    print("=" * 60)

    # 加载数据
    print("\n[1/3] Loading data...")
    data = load_data()
    print(f"  [OK] models: {len(data['model'])}")
    print(f"  [OK] causal ATT={data['causal']['att']['att']:.4f}")
    print(f"  [OK] monthly KPI: {len(data['kpi'])} months")
    print(f"  [OK] categories: {len(data['category'])}")

    # 创建PPT
    print("\n[2/3] Generating slides...")
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # 按顺序生成18页
    slide_cover(prs)
    print("  [OK] Slide 1: Cover")
    slide_toc(prs)
    print("  [OK] Slide 2: TOC")
    slide_background(prs)
    print("  [OK] Slide 3: Background")
    slide_objectives(prs)
    print("  [OK] Slide 4: Objectives")
    slide_architecture(prs)
    print("  [OK] Slide 5: Architecture")
    slide_data_overview(prs)
    print("  [OK] Slide 6: Data Overview")
    slide_rfm(prs)
    print("  [OK] Slide 7: RFM")
    slide_churn_definition(prs)
    print("  [OK] Slide 8: Churn Definition")
    slide_model_design(prs)
    print("  [OK] Slide 9: Model Design")
    slide_model_results(prs, data)
    print("  [OK] Slide 10: Model Results")
    slide_feature_importance(prs)
    print("  [OK] Slide 11: Feature Importance")
    slide_causal_design(prs)
    print("  [OK] Slide 12: Causal Design")
    slide_causal_results(prs, data)
    print("  [OK] Slide 13: Causal Results")
    slide_nlp(prs)
    print("  [OK] Slide 14: NLP")
    slide_business_findings(prs, data)
    print("  [OK] Slide 15: Business Findings")
    slide_recommendations(prs)
    print("  [OK] Slide 16: Recommendations")
    slide_highlights(prs)
    print("  [OK] Slide 17: Highlights")
    slide_qa(prs)
    print("  [OK] Slide 18: Q&A")

    # 保存
    print("\n[3/3] Saving...")
    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT_PATH))
    print(f"  [OK] Saved: {OUTPUT_PATH}")
    print(f"  [OK] Total slides: {len(prs.slides)}")
    print("\n" + "=" * 60)
    print("DONE!")
    print("=" * 60)


if __name__ == "__main__":
    main()
